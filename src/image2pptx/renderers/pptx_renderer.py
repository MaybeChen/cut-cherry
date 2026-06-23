from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from image2pptx.ir.elements import ElementType
from image2pptx.ir.slide_ir import SlideIR

EMU_PER_INCH=914400

def _hex_to_rgb(value: str | None) -> RGBColor:
    value=(value or "#ffffff").lstrip("#")
    return RGBColor(int(value[0:2],16), int(value[2:4],16), int(value[4:6],16))

class CoordinateMapper:
    def __init__(self, slide: SlideIR, prs: Presentation) -> None:
        self.sx = prs.slide_width / slide.width; self.sy = prs.slide_height / slide.height
    def box(self, x: float, y: float, w: float, h: float) -> tuple[Emu, Emu, Emu, Emu]:
        return Emu(int(x*self.sx)), Emu(int(y*self.sy)), Emu(int(w*self.sx)), Emu(int(h*self.sy))

class PptxRenderer:
    def render(self, ir: SlideIR, output_path: Path) -> Path:
        prs=Presentation(); prs.slide_width=Emu(EMU_PER_INCH*13.333); prs.slide_height=Emu(int(prs.slide_width*ir.height/ir.width))
        slide=prs.slides.add_slide(prs.slide_layouts[6]); mapper=CoordinateMapper(ir, prs)
        for e in ir.sort_by_z_index():
            x,y,w,h=mapper.box(e.bbox.x,e.bbox.y,e.bbox.width,e.bbox.height)
            if e.type == ElementType.BACKGROUND:
                shp=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x,y,w,h); shp.fill.solid(); shp.fill.fore_color.rgb=_hex_to_rgb(e.style.fill_color); shp.line.fill.background()
            elif e.type == ElementType.SHAPE:
                shp=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if e.style.shape_type=="roundRect" else MSO_SHAPE.RECTANGLE, x,y,w,h); shp.fill.solid(); shp.fill.fore_color.rgb=_hex_to_rgb(e.style.fill_color); shp.line.color.rgb=_hex_to_rgb(e.style.line_color)
            elif e.type == ElementType.TEXT:
                tb=slide.shapes.add_textbox(x,y,w,h); p=tb.text_frame.paragraphs[0]; p.text=e.text or ""; p.font.size=Pt(e.style.font_size or max(8, e.bbox.height*0.45)); p.font.color.rgb=_hex_to_rgb(e.style.font_color)
            elif e.type == ElementType.CONNECTOR:
                slide.shapes.add_connector(1, x, y, Emu(x+w), Emu(y+h))
            elif e.asset_path:
                slide.shapes.add_picture(str(e.asset_path), x,y,w,h)
        output_path.parent.mkdir(parents=True, exist_ok=True); prs.save(output_path); return output_path
